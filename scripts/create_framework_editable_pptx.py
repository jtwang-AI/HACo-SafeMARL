from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path(__file__).resolve().parents[1] / "submission" / "fig1_framework_editable.pptx"


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_textbox(slide, x, y, w, h, text, size=10, bold=False, color="#222222", align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
    return box


def add_card(slide, x, y, w, h, text, fill="#FFFFFF", line="#8C8C8C", size=10, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.85)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line_text in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line_text
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb("#151515")
    return shape


def add_rect(slide, x, y, w, h, fill="#FFFFFF", line="#999999", radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.85)
    return shape


def add_line(slide, x1, y1, x2, y2, color="#666666", width=1.1, arrow=True, dash=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width)
    if arrow:
        conn.line.end_arrowhead = True
    if dash:
        conn.line.dash_style = 4
    return conn


def add_polyline(slide, pts, color="#666666", width=1.1, arrow=True):
    for i in range(len(pts) - 1):
        add_line(slide, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1],
                 color=color, width=width, arrow=(arrow and i == len(pts) - 2))


def add_ship(slide, x, y, color="#1738D4", label="USV\nrelay"):
    hull = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Inches(x), Inches(y), Inches(0.55), Inches(0.18))
    hull.fill.solid()
    hull.fill.fore_color.rgb = rgb(color)
    hull.line.color.rgb = rgb(color)
    cabin = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.21), Inches(y - 0.16), Inches(0.18), Inches(0.16))
    cabin.fill.solid()
    cabin.fill.fore_color.rgb = rgb("#FFFFFF")
    cabin.line.color.rgb = rgb(color)
    add_textbox(slide, x - 0.05, y - 0.56, 0.65, 0.35, label, size=10, color=color)


def add_auv(slide, x, y, label):
    body = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.30), Inches(0.13))
    body.fill.solid()
    body.fill.fore_color.rgb = rgb("#13A620")
    body.line.color.rgb = rgb("#0D8218")
    tail = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x + 0.26), Inches(y + 0.005), Inches(0.16), Inches(0.12))
    tail.rotation = 90
    tail.fill.solid()
    tail.fill.fore_color.rgb = rgb("#13A620")
    tail.line.color.rgb = rgb("#0D8218")
    add_textbox(slide, x - 0.03, y + 0.13, 0.48, 0.20, label, size=9, color="#0B7A17")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # White background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb("#FFFFFF")

    # Stage containers
    stages = [
        (0.80, 1.26, 2.38, 4.04, "Mission and\nacoustic sensing"),
        (3.38, 1.26, 2.36, 4.04, "Heterogeneous\ngraph state"),
        (6.00, 1.26, 2.92, 4.04, "CTDE learning\nand optimization"),
        (9.78, 1.26, 2.38, 4.04, "Shielded\nexecution"),
    ]
    for x, y, w, h, title in stages:
        add_rect(slide, x, y, w, h, fill="#FAFAFA", line="#9E9E9E")
        add_textbox(slide, x + 0.12, y + 0.10, w - 0.24, 0.43, title, size=12, bold=True, color="#333333")

    # Mission scene
    add_rect(slide, 1.05, 1.88, 1.70, 1.28, fill="#E7E7FF", line="#E7E7FF", radius=False)
    add_rect(slide, 1.05, 3.16, 1.70, 1.23, fill="#D8D8FF", line="#D8D8FF", radius=False)
    for i in range(5):
        add_line(slide, 1.05 + i * 0.34, 3.16, 1.22 + i * 0.34, 3.12, color="#0B36B8", width=0.8, arrow=False)
    add_ship(slide, 1.43, 2.41, "#0B32D9", "USV\nrelay")
    add_ship(slide, 2.27, 2.45, "#EF3333", "traffic")
    for ax, ay, label in [(1.20, 3.55, "AUV$_1$"), (1.66, 3.90, "AUV$_2$"), (2.18, 3.55, "AUV$_3$")]:
        add_auv(slide, ax, ay, label.replace("$_", "").replace("$", ""))
        add_line(slide, 1.71, 2.59, ax + 0.15, ay + 0.04, color="#D97B00", width=0.9, arrow=False, dash=True)
    for ox, oy, sz in [(2.44, 4.05, 0.14), (2.52, 3.92, 0.10)]:
        obs = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ox), Inches(oy), Inches(sz), Inches(sz))
        obs.fill.solid()
        obs.fill.fore_color.rgb = rgb("#B8B8B8")
        obs.line.color.rgb = rgb("#858585")
    add_card(slide, 1.08, 4.36, 1.65, 0.63, "USBL/ray margins\npacket loss, delay", fill="#ECFFF0", size=10)

    # Graph state cards
    add_card(slide, 3.62, 1.80, 1.70, 0.70, "Typed nodes\nUSV, AUV,\nvessel, task", fill="#F0F0FF", size=10)
    add_card(slide, 3.62, 2.59, 1.70, 0.70, "Typed edges\nacoustic, collision,\nCOLREGs", fill="#F0F0FF", size=10)
    add_card(slide, 3.62, 3.49, 1.70, 0.54, "Graph H_t\nvariable team size", fill="#FFF7FF", size=10)
    add_card(slide, 3.62, 4.16, 1.70, 0.74, "Constraint signals\nc_out, collision\nenergy", fill="#ECFFF0", size=10)
    add_line(slide, 4.98, 2.50, 4.98, 2.59, arrow=True, width=0.8)
    add_line(slide, 4.98, 3.29, 4.98, 3.49, arrow=True, width=0.8)
    add_line(slide, 4.13, 4.16, 4.13, 4.03, arrow=True, width=0.8)

    # Learning cards
    add_card(slide, 6.32, 1.86, 2.00, 0.54, "GA--PSO--TLBO\n+ classical planners", fill="#FFF3E6", size=10)
    add_card(slide, 6.32, 2.62, 2.00, 0.54, "Type encoders\n+ graph attention", fill="#F0F0FF", size=10)
    add_card(slide, 6.32, 3.38, 2.00, 0.54, "Recurrent memory\nfor delayed links", fill="#F0F0FF", size=10)
    add_card(slide, 6.62, 4.05, 1.40, 0.54, "Actors\nUSV/AUV", fill="#FFF7FF", size=10)
    add_card(slide, 6.32, 4.78, 2.00, 0.54, "Reward/constraint\ncritics + multipliers", fill="#ECFFF0", size=10)
    add_line(slide, 7.65, 2.40, 7.65, 2.62, arrow=True, width=0.8)
    add_line(slide, 7.65, 3.16, 7.65, 3.38, arrow=True, width=0.8)
    add_line(slide, 7.52, 3.92, 7.52, 4.05, arrow=False, width=0.75)
    add_line(slide, 7.18, 4.59, 7.18, 4.78, arrow=False, width=0.75)

    # Execution cards
    add_card(slide, 10.02, 2.38, 1.86, 0.70, "Dual shields\nsurface/underwater\nCOLREGs, obstacles", fill="#FFF1F1", size=10)
    add_card(slide, 10.06, 3.45, 1.78, 0.54, "Safe actions\na_t -> a_t^safe", fill="#ECFFF0", size=10)
    add_card(slide, 10.06, 4.20, 1.78, 0.72, "Logged metrics\noutage, collisions\ntask progress", fill="#F0F0FF", size=10)
    add_line(slide, 11.02, 3.08, 11.02, 3.45, arrow=True, width=0.8)
    add_line(slide, 11.02, 3.99, 11.02, 4.20, arrow=True, width=0.8)

    # Cross-stage arrows
    add_line(slide, 2.74, 2.46, 3.62, 2.05, width=1.2)
    add_line(slide, 2.74, 3.05, 3.62, 2.94, width=1.2)
    add_line(slide, 2.73, 4.65, 3.62, 4.53, width=1.2)
    add_polyline(slide, [(5.32, 3.76), (5.88, 3.76), (5.88, 2.89), (6.32, 2.89)], width=1.2)
    add_polyline(slide, [(5.32, 4.53), (5.88, 4.53), (5.88, 5.05), (6.32, 5.05)], width=1.2)
    add_polyline(slide, [(8.02, 4.32), (9.45, 4.32), (9.45, 3.72), (10.06, 3.72)], width=1.2)

    # Small footer note for editability, outside the figure area.
    add_textbox(
        slide,
        0.78,
        6.78,
        11.4,
        0.22,
        "Editable recreation of Fig. 1 framework diagram: all boxes, labels, arrows, and simple scene elements are native PowerPoint objects.",
        size=7,
        color="#777777",
        align=PP_ALIGN.LEFT,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
